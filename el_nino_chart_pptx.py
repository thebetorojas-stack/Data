"""
El Nino exposure vs. fiscal fragility - as a NATIVE, EDITABLE PowerPoint slide.
Every bubble, label and quadrant is a real PowerPoint shape, so you can drag,
recolor or retype anything directly in PowerPoint after running this.

Needs only python-pptx:  pip install python-pptx
Run it -> writes 'el_nino_exposure_chart.pptx'.

EDIT ONLY THE `DATA` DICT.
  x = El Nino physical exposure (0 = insulated, 10 = severe)
  y = fiscal / macro fragility  (0 = resilient, 10 = fragile)
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.oxml.ns import qn

# ---------------------------------------------------------------------------
# DATA  (exposure, fragility) — adjust freely
# ---------------------------------------------------------------------------
DATA = {
    "Colombia":  (8.0, 8.5),
    "Peru":      (9.0, 5.0),
    "Brazil":    (6.2, 8.0),
    "Argentina": (2.2, 8.5),   # El Nino HELPS crops; fragile = adjustment risk
    "Panama":    (4.5, 6.0),
    "Venezuela": (5.5, 9.0),   # Guri hydro exposed to drought
    "Mexico":    (3.5, 6.0),
    "Dom. Rep.": (4.4, 4.8),
    "Chile":     (3.0, 3.0),
}
HIGHLIGHT = {"Colombia", "Peru", "Brazil"}   # drawn larger

# colours
ACCENT = RGBColor(0xC8, 0x10, 0x2E)   # UBS-ish red
GREEN  = RGBColor(0x2E, 0x8B, 0x57)
GREY   = RGBColor(0x55, 0x55, 0x55)
DARK   = RGBColor(0x22, 0x22, 0x22)
LIGHTG = RGBColor(0xBB, 0xBB, 0xBB)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

# risk -> colour ramp (green -> amber -> red), exposure-weighted
_STOPS = [(0.0, (0x2E, 0x8B, 0x57)), (0.5, (0xE6, 0xB8, 0x00)), (1.0, (0xC8, 0x10, 0x2E))]
def risk_color(e, f):
    t = (0.62 * e + 0.38 * f - 3.0) / (8.5 - 3.0)   # normalise to ~0..1
    t = max(0.0, min(1.0, t))
    for i in range(len(_STOPS) - 1):
        t0, c0 = _STOPS[i]; t1, c1 = _STOPS[i + 1]
        if t <= t1:
            k = (t - t0) / (t1 - t0)
            return RGBColor(*[int(c0[j] + (c1[j] - c0[j]) * k) for j in range(3)])
    return RGBColor(*_STOPS[-1][1])

# ---------------------------------------------------------------------------
# plot geometry (inches) inside a 13.333 x 7.5 slide
# ---------------------------------------------------------------------------
PL, PR = 1.55, 12.45      # plot left / right
PT, PB = 0.95, 6.45       # plot top / bottom
def px(x): return Inches(PL + (PR - PL) * (x / 10.0))
def py(y): return Inches(PB - (PB - PT) * (y / 10.0))   # invert

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
slide = prs.slides.add_slide(prs.slide_layouts[6])   # blank

def no_line(shape):
    shape.line.fill.background()

def add_rect(x0, y0, x1, y1, rgb, alpha=None):
    left, top = px(x0), py(y1)
    w = px(x1) - px(x0); h = py(y0) - py(y1)
    sp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, h)
    sp.fill.solid(); sp.fill.fore_color.rgb = rgb
    no_line(sp); sp.shadow.inherit = False
    if alpha is not None:   # set fill transparency via XML (0-100000 = %)
        srgb = sp.fill.fore_color._xFill.find(qn('a:srgbClr'))
        a = srgb.makeelement(qn('a:alpha'), {'val': str(int(alpha * 1000))})
        srgb.append(a)
    return sp

def add_dashed_line(x0, y0, x1, y1):
    ln = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, px(x0), py(y0), px(x1), py(y1))
    ln.line.color.rgb = LIGHTG; ln.line.width = Pt(1)
    d = ln.line._get_or_add_ln().makeelement(qn('a:prstDash'), {'val': 'dash'})
    ln.line._get_or_add_ln().append(d)
    return ln

def add_text(x_in, y_in, w_in, h_in, text, size, color, bold=False,
             align=PP_ALIGN.CENTER, italic=False, rot=0, anchor=MSO_ANCHOR.MIDDLE):
    tb = slide.shapes.add_textbox(Inches(x_in), Inches(y_in), Inches(w_in), Inches(h_in))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text
    r.font.size = Pt(size); r.font.bold = bold; r.font.italic = italic
    r.font.color.rgb = color; r.font.name = "Arial"
    if rot: tb.rotation = rot
    return tb

# ---- quadrant tints -------------------------------------------------------
add_rect(5, 5, 10, 10, ACCENT, alpha=6)    # top-right: danger
add_rect(0, 0, 5, 5, GREEN, alpha=6)       # bottom-left: safe

# ---- midlines -------------------------------------------------------------
add_dashed_line(5, 0, 5, 10)
add_dashed_line(0, 5, 10, 5)

# ---- quadrant corner labels ----------------------------------------------
add_text(PR - 3.3, PT + 0.02, 3.2, 0.35, "MOST VULNERABLE", 13, ACCENT, bold=True, align=PP_ALIGN.RIGHT)
add_text(PL + 0.05, PT + 0.02, 3.8, 0.35, "Fragile, but off El Niño's path", 11, GREY, italic=True, align=PP_ALIGN.LEFT)
add_text(PL + 0.05, PB - 0.37, 3.0, 0.35, "MOST INSULATED", 13, GREEN, bold=True, align=PP_ALIGN.LEFT)
add_text(PR - 3.3, PB - 0.37, 3.2, 0.35, "Exposed, but able to absorb", 11, GREY, italic=True, align=PP_ALIGN.RIGHT)

# ---- bubbles + country labels --------------------------------------------
for name, (e, f) in DATA.items():
    dia = 0.62 if name in HIGHLIGHT else 0.50
    cx, cy = px(e), py(f)
    sp = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                cx - Emu(int(Inches(dia) / 2)), cy - Emu(int(Inches(dia) / 2)),
                                Inches(dia), Inches(dia))
    sp.fill.solid(); sp.fill.fore_color.rgb = risk_color(e, f)
    sp.line.color.rgb = WHITE; sp.line.width = Pt(1.5); sp.shadow.inherit = False
    # label just under the bubble
    add_text(cx / 914400 - 1.0, cy / 914400 + dia / 2 + 0.02, 2.0, 0.3,
             name, 11, DARK, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

# ---- axes labels ----------------------------------------------------------
add_text(PL - 0.05, PB + 0.05, 1.2, 0.3, "Low", 11, GREY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP)
add_text(PR - 1.15, PB + 0.05, 1.2, 0.3, "High", 11, GREY, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.TOP)
add_text(PL + (PR - PL) / 2 - 2.0, PB + 0.42, 4.0, 0.35,
         "El Niño physical exposure  →", 12.5, DARK, align=PP_ALIGN.CENTER)
# y-axis: rotated text boxes
add_text(PL - 1.05, PT - 0.02, 1.2, 0.3, "Fragile", 11, GREY, rot=270, align=PP_ALIGN.CENTER)
add_text(PL - 1.05, PB - 0.28, 1.2, 0.3, "Resilient", 11, GREY, rot=270, align=PP_ALIGN.CENTER)
add_text(PL - 1.35, PT + (PB - PT) / 2 - 1.5, 3.0, 0.35,
         "Fiscal / macro fragility  →", 12.5, DARK, rot=270, align=PP_ALIGN.CENTER)

prs.save("el_nino_exposure_chart.pptx")
print("saved el_nino_exposure_chart.pptx")
